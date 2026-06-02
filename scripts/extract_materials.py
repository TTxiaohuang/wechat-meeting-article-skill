#!/usr/bin/env python3
"""Extract text from meeting material folders into agent-friendly Markdown."""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any


DEPENDENCY_HELP = """
Missing optional document parser dependencies.

Install them in the active Python environment:
  python -m pip install python-docx python-pptx pdfplumber pypdf

For .doc files (old Word format), the script auto-converts to .docx when possible.
Requires either Microsoft Word (via pywin32) or LibreOffice on PATH.

On Windows, if Chinese output is garbled, run:
  $env:PYTHONIOENCODING="utf-8"
"""


def configure_stdio() -> None:
    for stream in (sys.stdout, sys.stderr):
        if hasattr(stream, "reconfigure"):
            stream.reconfigure(encoding="utf-8")


def slug(value: str) -> str:
    value = re.sub(r'["“”‘’`´]+', "_", value)
    value = value.replace("：", "_").replace("、", "_")
    cleaned = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', "_", value)
    cleaned = re.sub(r"_+", "_", cleaned).strip(" ._")
    return cleaned[:90] or "material"


def read_text_file(path: Path) -> str:
    data = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "gbk"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def convert_doc_to_docx(doc_path: Path) -> Path:
    """Convert a .doc file to .docx. Returns the path to the temporary .docx file."""
    import shutil
    import subprocess
    import tempfile

    tmp_dir = Path(tempfile.mkdtemp(prefix="doc_convert_"))
    docx_path = tmp_dir / (doc_path.stem + ".docx")

    # Try Word COM automation (Windows + pywin32)
    try:
        import win32com.client  # type: ignore[import-untyped]
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        try:
            doc = word.Documents.Open(str(doc_path.resolve()))
            doc.SaveAs2(str(docx_path.resolve()), FileFormat=16)  # 16 = docx
            doc.Close()
        finally:
            word.Quit()
        return docx_path
    except Exception:
        pass

    # Try LibreOffice headless
    lo_cmd = shutil.which("libreoffice") or shutil.which("soffice")
    if lo_cmd:
        try:
            subprocess.run(
                [lo_cmd, "--headless", "--convert-to", "docx", "--outdir", str(tmp_dir), str(doc_path.resolve())],
                check=True, capture_output=True, timeout=60,
            )
            if docx_path.is_file():
                return docx_path
        except Exception:
            pass

    raise RuntimeError(
        f"Cannot convert .doc file: {doc_path.name}. "
        "Install Microsoft Word + pywin32, or install LibreOffice and add it to PATH. "
        "Alternatively, manually save the file as .docx first."
    )


def extract_html(path: Path) -> str:
    """Extract readable text from an HTML file, stripping scripts and styles."""
    from html.parser import HTMLParser

    raw = read_text_file(path)

    # Quick extraction for reveal.js / Slidev / similar slide frameworks
    # These use <section> or <div class="slide"> to wrap each slide
    slide_pattern = re.compile(
        r'<(?:section|div)\b[^>]*class="[^"]*(?:slide|reveal)[^"]*"[^>]*>(.*?)</(?:section|div)>',
        re.DOTALL | re.IGNORECASE,
    )
    slides = slide_pattern.findall(raw)

    class _TextExtractor(HTMLParser):
        def __init__(self) -> None:
            super().__init__()
            self._skip = False
            self._skip_tags = {"script", "style", "noscript", "svg", "head"}
            self._parts: list[str] = []

        def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
            if tag in self._skip_tags:
                self._skip = True
            if tag in {"h1", "h2", "h3", "h4", "h5", "h6"}:
                self._parts.append("\n\n## ")
            elif tag in {"p", "div", "br"}:
                self._parts.append("\n\n")
            elif tag == "li":
                self._parts.append("\n- ")
            elif tag == "tr":
                self._parts.append("\n")
            elif tag == "td" or tag == "th":
                self._parts.append(" | ")
            elif tag == "hr":
                self._parts.append("\n---\n")

        def handle_endtag(self, tag: str) -> None:
            if tag in self._skip_tags:
                self._skip = False
            if tag in {"h1", "h2", "h3", "h4", "h5", "h6"}:
                self._parts.append("\n")

        def handle_data(self, data: str) -> None:
            if not self._skip:
                self._parts.append(data)

    def _extract_block(block: str) -> str:
        parser = _TextExtractor()
        parser.feed(block)
        text = "".join(parser._parts)
        text = re.sub(r'[ \t]+', ' ', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

    # If slide sections found, extract each as a separate block
    if len(slides) >= 2:
        results: list[str] = []
        for i, slide_html in enumerate(slides, start=1):
            text = _extract_block(slide_html)
            if text:
                results.append(f"## Slide {i}\n\n{text}")
        if results:
            return "\n\n".join(results)

    # General HTML extraction
    return _extract_block(raw)
    """Extract plain text from an RTF file by stripping control codes."""
    raw = read_text_file(path)
    # Remove header up to first \deflang or \ansi or first group content
    text = raw
    # Strip {\*\...} destination groups
    text = re.sub(r'\{\\\*[^{}]*\}', '', text)
    # Strip control words (e.g. \rtf1, \ansi, \fonttbl, \colortbl, \stylesheet)
    text = re.sub(r'\\[a-zA-Z]+\d* ?(?=\{|\}|\\)', '', text)
    # Strip escaped characters like \'xx
    text = re.sub(r"\\'[0-9a-fA-F]{2}", '', text)
    # Strip remaining backslash commands like \\, \{, \}
    text = re.sub(r'\\[{}\\]', '', text)
    # Remove braces
    text = text.replace('{', '').replace('}', '')
    # Collapse whitespace
    text = re.sub(r'[ \t]+', ' ', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    # Convert RTF unicode escapes \uN? to approximate char
    text = re.sub(r'\\u(\-?\d+)\??', lambda m: chr(int(m.group(1)) % 65536), text)
    return text.strip()


def extract_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is required for .docx files") from exc

    doc = Document(str(path))
    parts: list[str] = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for .pptx files") from exc

    prs = Presentation(str(path))
    slides: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    texts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    if any(cells):
                        texts.append(" | ".join(cells))
        if texts:
            slides.append(f"## Slide {index}\n\n" + "\n\n".join(texts))
    return "\n\n".join(slides)


def parse_pdf_pages(value: str) -> int | None:
    normalized = str(value).strip().lower()
    if normalized in {"all", "full", "*"}:
        return None
    try:
        pages = int(normalized)
    except ValueError as exc:
        raise argparse.ArgumentTypeError("--pdf-pages must be a positive integer or 'all'") from exc
    if pages <= 0:
        raise argparse.ArgumentTypeError("--pdf-pages must be a positive integer or 'all'")
    return pages


def requested_pages_label(value: int | None) -> str:
    return "all" if value is None else str(value)


def extract_pdf(path: Path, max_pages: int | None) -> tuple[str, dict[str, Any]]:
    try:
        import pdfplumber
    except ImportError:
        pdfplumber = None

    if pdfplumber is not None:
        pages: list[str] = []
        total_pages = 0
        with pdfplumber.open(str(path)) as pdf:
            total_pages = len(pdf.pages)
            selected_pages = pdf.pages if max_pages is None else pdf.pages[:max_pages]
            for index, page in enumerate(selected_pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"## Page {index}\n\n{text.strip()}")
        metadata = {
            "pdf_pages_total": total_pages,
            "pdf_pages_extracted": total_pages if max_pages is None else min(max_pages, total_pages),
            "pdf_pages_requested": requested_pages_label(max_pages),
            "pdf_text_pages": len(pages),
        }
        if pages:
            return "\n\n".join(pages), metadata

    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("pdfplumber or pypdf is required for .pdf files") from exc

    reader = PdfReader(str(path))
    pages = []
    total_pages = len(reader.pages)
    selected_pages = reader.pages if max_pages is None else reader.pages[:max_pages]
    for index, page in enumerate(selected_pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"## Page {index}\n\n{text.strip()}")
    return "\n\n".join(pages), {
        "pdf_pages_total": total_pages,
        "pdf_pages_extracted": total_pages if max_pages is None else min(max_pages, total_pages),
        "pdf_pages_requested": requested_pages_label(max_pages),
        "pdf_text_pages": len(pages),
    }


def extract_one(path: Path, pdf_pages: int | None) -> tuple[str, str, dict[str, Any]]:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return "docx", extract_docx(path), {}
    if suffix == ".doc":
        docx_path = convert_doc_to_docx(path)
        return "doc", extract_docx(docx_path), {"converted_from": ".doc"}
    if suffix == ".pptx":
        return "pptx", extract_pptx(path), {}
    if suffix == ".pdf":
        text, metadata = extract_pdf(path, pdf_pages)
        return "pdf", text, metadata
    if suffix == ".rtf":
        return "rtf", extract_rtf(path), {}
    if suffix in {".txt", ".md"}:
        return suffix.lstrip("."), read_text_file(path), {}
    if suffix in {".html", ".htm"}:
        return "html", extract_html(path), {}
    raise ValueError(f"unsupported file type: {suffix}")


def main() -> int:
    configure_stdio()
    parser = argparse.ArgumentParser()
    parser.add_argument("input_dir", type=Path)
    parser.add_argument("--out", type=Path, default=Path("extracted_materials"))
    parser.add_argument("--pdf-pages", type=parse_pdf_pages, default=None)
    parser.add_argument("--no-recursive", action="store_true",
                        help="Only scan top-level files, do not recurse into subdirectories")
    args = parser.parse_args()

    input_dir = args.input_dir.resolve()
    out_dir = args.out.resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    supported = {".docx", ".doc", ".pptx", ".pdf", ".txt", ".md", ".rtf", ".html", ".htm"}
    if args.no_recursive:
        files = [p for p in sorted(input_dir.iterdir()) if p.is_file() and p.suffix.lower() in supported]
    else:
        files = [p for p in sorted(input_dir.rglob("*")) if p.is_file() and p.suffix.lower() in supported]
    manifest: dict[str, Any] = {"input_dir": str(input_dir), "files": []}
    had_error = False

    for path in files:
        rel = path.relative_to(input_dir)
        if len(rel.parts) > 1:
            prefix = slug("_".join(rel.parts[:-1]))
            output_name = f"{prefix}_{slug(path.stem)}.md"
        else:
            output_name = f"{slug(path.stem)}.md"
        output_path = out_dir / output_name
        record: dict[str, Any] = {
            "source": str(path),
            "type": path.suffix.lower().lstrip("."),
            "output": str(output_path),
            "status": "ok",
        }
        try:
            detected_type, text, metadata = extract_one(path, args.pdf_pages)
            record["type"] = detected_type
            record.update(metadata)
            if detected_type == "pdf" and metadata.get("pdf_text_pages") == 0:
                record["warning"] = "No extractable text found in selected PDF pages; the PDF may be scanned or image-based."
            output_path.write_text(
                f"# {path.name}\n\nSource: `{path}`\n\n{text.strip()}\n",
                encoding="utf-8",
            )
            record["chars"] = len(text)
        except Exception as exc:  # noqa: BLE001 - report per-file extraction failures.
            had_error = True
            record["status"] = "error"
            record["error"] = str(exc)
        manifest["files"].append(record)

    (out_dir / "materials_manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    scope = "top-level" if args.no_recursive else "recursive"
    print(f"Scanned {len(files)} supported files from {input_dir} ({scope})")
    print(f"Wrote manifest: {out_dir / 'materials_manifest.json'}")
    if had_error:
        print(DEPENDENCY_HELP, file=sys.stderr)
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
