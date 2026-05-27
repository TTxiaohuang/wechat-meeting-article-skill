#!/usr/bin/env python3
"""Extract text from meeting material folders into agent-friendly Markdown."""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any


DEPENDENCY_HELP = """
Missing optional document parser dependencies.

Install them in the active Python environment:
  python -m pip install python-docx python-pptx pdfplumber pypdf

On Windows, if Chinese output is garbled, run:
  $env:PYTHONIOENCODING="utf-8"
"""


def configure_stdio() -> None:
    for stream in (sys.stdout, sys.stderr):
        if hasattr(stream, "reconfigure"):
            stream.reconfigure(encoding="utf-8")


def slug(value: str) -> str:
    value = re.sub(r'["“”‘’`´]+', "_", value)
    value = value.replace("：", "_").replace("、", "_")
    cleaned = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', "_", value)
    cleaned = re.sub(r"_+", "_", cleaned).strip(" ._")
    return cleaned[:90] or "material"


def read_text_file(path: Path) -> str:
    data = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "gbk"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def extract_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is required for .docx files") from exc

    doc = Document(str(path))
    parts: list[str] = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for .pptx files") from exc

    prs = Presentation(str(path))
    slides: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    texts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    if any(cells):
                        texts.append(" | ".join(cells))
        if texts:
            slides.append(f"## Slide {index}\n\n" + "\n\n".join(texts))
    return "\n\n".join(slides)


def extract_pdf(path: Path, max_pages: int) -> str:
    try:
        import pdfplumber
    except ImportError:
        pdfplumber = None

    if pdfplumber is not None:
        pages: list[str] = []
        with pdfplumber.open(str(path)) as pdf:
            for index, page in enumerate(pdf.pages[:max_pages], start=1):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"## Page {index}\n\n{text.strip()}")
        if pages:
            return "\n\n".join(pages)

    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("pdfplumber or pypdf is required for .pdf files") from exc

    reader = PdfReader(str(path))
    pages = []
    for index, page in enumerate(reader.pages[:max_pages], start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"## Page {index}\n\n{text.strip()}")
    return "\n\n".join(pages)


def extract_one(path: Path, pdf_pages: int) -> tuple[str, str]:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return "docx", extract_docx(path)
    if suffix == ".pptx":
        return "pptx", extract_pptx(path)
    if suffix == ".pdf":
        return "pdf", extract_pdf(path, pdf_pages)
    if suffix in {".txt", ".md"}:
        return suffix.lstrip("."), read_text_file(path)
    raise ValueError(f"unsupported file type: {suffix}")


def main() -> int:
    configure_stdio()
    parser = argparse.ArgumentParser()
    parser.add_argument("input_dir", type=Path)
    parser.add_argument("--out", type=Path, default=Path("extracted_materials"))
    parser.add_argument("--pdf-pages", type=int, default=5)
    args = parser.parse_args()

    input_dir = args.input_dir.resolve()
    out_dir = args.out.resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    supported = {".docx", ".pptx", ".pdf", ".txt", ".md"}
    files = [p for p in sorted(input_dir.iterdir()) if p.is_file() and p.suffix.lower() in supported]
    manifest: dict[str, Any] = {"input_dir": str(input_dir), "files": []}
    had_error = False

    for path in files:
        output_name = f"{slug(path.stem)}.md"
        output_path = out_dir / output_name
        record: dict[str, Any] = {
            "source": str(path),
            "type": path.suffix.lower().lstrip("."),
            "output": str(output_path),
            "status": "ok",
        }
        try:
            detected_type, text = extract_one(path, args.pdf_pages)
            record["type"] = detected_type
            output_path.write_text(
                f"# {path.name}\n\nSource: `{path}`\n\n{text.strip()}\n",
                encoding="utf-8",
            )
            record["chars"] = len(text)
        except Exception as exc:  # noqa: BLE001 - report per-file extraction failures.
            had_error = True
            record["status"] = "error"
            record["error"] = str(exc)
        manifest["files"].append(record)

    (out_dir / "materials_manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    print(f"Scanned {len(files)} supported files from {input_dir}")
    print(f"Wrote manifest: {out_dir / 'materials_manifest.json'}")
    if had_error:
        print(DEPENDENCY_HELP, file=sys.stderr)
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
